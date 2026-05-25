import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    first = True
    for bullet in content:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        p.text = bullet
        p.font.size = Pt(20)
        p.font.name = "Calibri"
    return slide

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Tying the Data Knot"
    subtitle.text = "Predicting Meaningful Connections\nGroup Assignment (WIA1006/WID3006)\nSem 2, Session 2025/2026 | FCSIT, Universiti Malaya"

    # Slide 1: Project Overview
    add_slide(prs, "Project Overview", [
        "Objective: Predict whether a dating app user will achieve a meaningful connection.",
        "Based on demographic profile and in-app behaviour patterns.",
        "ML Task Type: Binary Classification",
        " - Positive: Mutual Match, Instant Match, Date Happened, Relationship Formed",
        " - Negative: Ghosted, Blocked, Catfished, Chat Ignored, No Action, One-sided Like",
        "Dataset: dating_app_behavior_dataset_extended1.csv (50,000 records × 25 features)"
    ])

    # Slide 2: Why We Use the Extended Dataset
    add_slide(prs, "The Extended Dataset", [
        "Adds 6 new features not present in the original dataset:",
        "1. Age (18–59): Core dating preference factor",
        "2. Height (145–200 cm) & 3. Weight: Physical profile signals",
        "4. Body Type (6 types): Profile completeness & preference signal",
        "5. Relationship Intent (6 types): Strong predictor (e.g. Serious vs Hookups)",
        "6. Zodiac Sign (12 signs): Cultural/personality correlation",
        "Dataset Quality: Zero missing values, zero duplicates, perfectly balanced target."
    ])

    # Slide 3: Data Preprocessing
    add_slide(prs, "Data Preprocessing Pipeline", [
        "1. Dropped redundant columns (app_usage_time_label, swipe_right_label).",
        "2. Binary Target creation: 10 classes -> Positive (39.7%) / Negative (60.3%).",
        "3. Ordinal Encoding: income_bracket, education_level (3 tiers).",
        "4. One-Hot Encoding: 7 nominal columns -> 43 binary columns.",
        "5. Multi-Hot Encoding: interest_tags (49 unique tags).",
        "6. StandardScaler Normalization: Applied to all 12 numeric features.",
        "Result: 50,000 rows × 113 feature columns + 1 target column."
    ])

    # Slide 4: Feature Selection
    add_slide(prs, "Feature Selection & PCA", [
        "Goal: Reduce 113 features to a more informative, smaller set.",
        "Method 1: ANOVA F-Score (SelectKBest) - Evaluates statistical differences between classes.",
        "Method 2: Mutual Information - Captures non-linear relationships.",
        "Final Selection: Union strategy yielded 67 most important features.",
        "PCA (Dimensionality Reduction):",
        " - Applied PCA to the 67 selected features.",
        " - Maintained 95.2% of variance with 55 components.",
        "Two feature sets generated: X_selected (67 features) and X_pca (55 components)."
    ])
    
    # Slide 5: Model Training
    add_slide(prs, "Model Training & Evaluation", [
        "6 models trained on an 80/20 stratified split:",
        " - Logistic Regression, K-Nearest Neighbors (KNN), Decision Tree,",
        " - Random Forest, XGBoost, Support Vector Machine (SVM).",
        "Evaluation Metrics: Accuracy, Precision, Recall, F1 Score, ROC-AUC.",
        "Hardware Optimizations: Accelerated 16-Thread SVM Bagging Ensemble to reduce training time.",
        "Cross-Validation Parallel Optimization to utilize CPU effectively."
    ])

    # Slide 6: Hyperparameter Tuning
    add_slide(prs, "Hyperparameter Tuning", [
        "Goal: Improve top 3 performing models.",
        "Method: RandomizedSearchCV with 30 combinations & 5-fold CV.",
        "Evaluated models:",
        " - Random Forest: Tuned n_estimators, max_depth, min_samples.",
        " - XGBoost: Tuned learning_rate, subsample, max_depth.",
        " - SVM: Tuned C, gamma, kernel.",
        "Output: Best parameters selected based on highest F1 score."
    ])

    # Slide 7: Advanced ML Enhancements
    add_slide(prs, "Advanced ML Techniques Implemented", [
        "1. Class Imbalance Mitigation: class_weight='balanced' prevents trivial majority-class predictors.",
        "2. Statistical Significance Testing: Paired t-test proved model performance gaps are significant (p=0.0004).",
        "3. SHAP Explainability: Added beeswarm plot for transparent decision-making visualization.",
        "4. Ethical Considerations: Analyzed demographic parity (accuracy across gender identities) to highlight algorithmic bias risks."
    ])

    # Slide 8: Conclusion
    add_slide(prs, "Conclusion", [
        "Features contained limited predictive signal (ROC-AUC ~0.50).",
        "Synthetic nature of data meant traditional relationships were absent.",
        "Demonstrated a fully robust, optimized ML pipeline from data loading to advanced evaluation.",
        "Utilized smart checkpointing and parallel processing for maximum efficiency.",
        "Pipeline ready for real-world datasets with actual underlying patterns."
    ])

    prs.save('reports/Tying_the_Data_Knot_Presentation.pptx')
    print("Presentation saved as reports/Tying_the_Data_Knot_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
