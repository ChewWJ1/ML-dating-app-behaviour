import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation(r'c:\Users\HP\Documents\GitHub\ML-dating-app-behaviour\reports\Tying the Data Knot Advanced Predictive Modeling for Connections.pptx')

for i, slide in enumerate(prs.slides):
    print("=== SLIDE", i+1, "===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(text)
    print()
